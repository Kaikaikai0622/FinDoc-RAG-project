"""PPTX 解析模块

使用 python-pptx 解析 .pptx 文件。
每张 slide 合并所有文本框内容 → category="Slide" 元素。
slide_title 取第一个文本框内容（截取前 50 字符）供 chunker 上下文注入。
"""
import logging
from pathlib import Path
from typing import List

from .models import ParsedDocument, ParsedElement

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]
    Pt = None  # type: ignore[assignment]
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 库不可用，无法解析 .pptx 文件。请执行: pip install python-pptx")


class PptxParser:
    """PPTX 解析器"""

    def parse(self, file_path: str) -> ParsedDocument:
        """解析 .pptx 文件

        Args:
            file_path: .pptx 文件路径

        Returns:
            ParsedDocument 对象

        Raises:
            RuntimeError: python-pptx 不可用
        """
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx 库不可用，无法解析 .pptx 文件。请执行: pip install python-pptx"
            )

        path = Path(file_path)
        source_file = path.name
        elements = self._parse_with_python_pptx(file_path)

        return ParsedDocument(
            source_file=source_file,
            file_type="pptx",
            elements=elements,
        )

    # ------------------------------------------------------------------
    # 内部方法
    # ------------------------------------------------------------------

    def _parse_with_python_pptx(self, file_path: str) -> List[ParsedElement]:
        """遍历每张 slide，合并文本框内容生成 ParsedElement"""
        elements: List[ParsedElement] = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            shape_texts: List[str] = []
            first_text = ""

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                shape_text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )
                if shape_text.strip():
                    if not first_text:
                        first_text = shape_text.strip()
                    shape_texts.append(shape_text)

            combined_text = "\n\n".join(shape_texts).strip()
            if not combined_text:
                continue

            # slide_title 取第一个文本框内容（最多50字符）
            slide_title = first_text[:50] if first_text else ""

            elements.append(ParsedElement(
                text=combined_text,
                category="Slide",
                page_or_index=slide_num,
                metadata={
                    "slide_number": slide_num,
                    "slide_title": slide_title,
                },
            ))

        return elements


def parse_pptx(file_path: str) -> ParsedDocument:
    """解析 PPTX 文件的便捷函数"""
    parser = PptxParser()
    return parser.parse(file_path)


